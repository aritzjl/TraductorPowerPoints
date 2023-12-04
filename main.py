import tkinter as tk
from tkinter import ttk
from tkinter import filedialog, messagebox
from pptx import Presentation
from deep_translator import GoogleTranslator

# Diccionario de códigos de idioma y nombres completos
idiomas = {
    "es": "Español",
    "uk": "Ucraniano",
    "en": "Inglés",
    "eu": "Euskera"
}

def translate_pptx():
    file_path = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
    if file_path:
        selected_language = combo_language.get()
        selected_language = next(key for key, value in idiomas.items() if value ==  selected_language)
        if selected_language in idiomas:
            target_language = idiomas[selected_language]
            output_file_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint files", "*.pptx")])
            if output_file_path:
                translate_button.config(state=tk.DISABLED)
                progress_bar.start()  # Inicia la barra de progreso al comenzar la traducción
                messagebox.showinfo("Inicio de traducción", "La traducción ha comenzado.")
                translate_pptx_file(file_path, output_file_path, selected_language, progress_bar)
                messagebox.showinfo("Fin de traducción", f"La traducción ha finalizado. Archivo guardado en: {output_file_path}")
                progress_bar.stop()  # Detiene la barra de progreso al finalizar la traducción
                translate_button.config(state=tk.NORMAL)

        else:
            messagebox.showwarning("Error", "Por favor selecciona un idioma válido.")

def translate_pptx_file(input_pptx, output_pptx, target_language, progress):
    translator = GoogleTranslator(source='auto', target=target_language)
    prs = Presentation(input_pptx)
    total_slides = len(prs.slides)
    translated_slides = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                translated_text = translator.translate(shape.text)
                shape.text = translated_text
        translated_slides += 1
        progress['value'] = (translated_slides / total_slides) * 100  # Actualiza el valor de la barra de progreso
        root.update_idletasks()  # Actualiza la interfaz gráfica para mostrar el progreso
    prs.save(output_pptx)

root = tk.Tk()
root.title("Traductor de PowerPoint")
style = ttk.Style(root)
root.tk.call("source", "forest-dark.tcl")
style.theme_use("forest-dark")
root.geometry("500x320")

label_titulo = ttk.Label(root, text="Traductor PowerPoints", font=("Helvetica", 16, "bold"))
label_titulo.pack(pady=10)
ttk.Separator(root, orient="horizontal", style="Accent.TButton").pack(fill="x")
label_file = ttk.Label(root, text="Selecciona un archivo PowerPoint:")
label_file.pack(pady=5)

button_browse = ttk.Button(root, text="Seleccionar archivo", command=translate_pptx)
button_browse.pack(pady=5)

label_language = ttk.Label(root, text="Selecciona el idioma de destino:")
label_language.pack(pady=5)

combo_language = ttk.Combobox(root, values=list(idiomas.values()))
combo_language.pack(pady=5)
combo_language.set("Ucraniano")  # Establece el valor predeterminado

translate_button = ttk.Button(root, text="Traducir", command=translate_pptx, state=tk.NORMAL, style="Accent.TButton")
translate_button.pack(pady=10)

progress_bar = ttk.Progressbar(root, orient=tk.HORIZONTAL, length=300, mode='determinate')
progress_bar.pack(pady=10)

watermark = ttk.Label(root, text="By AJL")
watermark.pack(side="right", padx=5)

root.mainloop()
